from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()

# Clean, professional light theme colors
bg_color = RGBColor(245, 247, 250) # Light slate
accent_blue = RGBColor(30, 60, 114) # Deep corporate blue
accent_red = RGBColor(220, 53, 69) # Professional red
text_dark = RGBColor(40, 45, 50)
text_light = RGBColor(100, 105, 110)

def set_light_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    
    # Add a top blue accent bar
    top_bar = slide.shapes.add_shape(
        1, 0, 0, Inches(10), Inches(0.2) # 1 is MSO_SHAPE.RECTANGLE
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = accent_blue
    top_bar.line.fill.background()

# 1. Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
set_light_bg(slide)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Delhivery Graph Intelligence"
title.text_frame.paragraphs[0].font.color.rgb = accent_blue
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.size = Pt(48)
subtitle.text = "Predicting Delays & Identifying Structural Bottlenecks\n\nPrepared by Shikhar Verma"
subtitle.text_frame.paragraphs[0].font.color.rgb = text_light
subtitle.text_frame.paragraphs[0].font.size = Pt(22)

# Add truck doodle to title slide
img_path = "outputs/figures/doodle_truck.png"
slide.shapes.add_picture(img_path, Inches(3.5), Inches(4.5), height=Inches(2.5))

def add_title(slide, text):
    title = slide.shapes.title
    title.text = text
    title.text_frame.paragraphs[0].font.color.rgb = accent_blue
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Pt(36)

def style_content(content):
    for p in content.text_frame.paragraphs:
        p.font.color.rgb = text_dark
        p.font.size = Pt(20)
        p.space_after = Pt(12)

# 2. The Problem
slide = prs.slides.add_slide(prs.slide_layouts[5])
set_light_bg(slide)
add_title(slide, "The Strategic Challenge")
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5), Inches(4))
content = txBox.text_frame
content.text = ("• Delhivery uses OSRM (shortest path) to estimate delivery times.\n"
               "• Real-world logistics is messy: dwell times, congestion, and operational constraints cause massive deviations.\n\n"
               "The Result:\n"
               "A 95%+ SLA Breach Rate on chronic routes, breaking downstream capacity planning.")
style_content(txBox)

img_path = "outputs/figures/doodle_bottleneck.png"
slide.shapes.add_picture(img_path, Inches(5.5), Inches(1.5), width=Inches(4))


# 3. The Solution
slide = prs.slides.add_slide(prs.slide_layouts[5])
set_light_bg(slide)
add_title(slide, "Our Graph-Based Solution")
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5), Inches(4))
content = txBox.text_frame
content.text = ("We modeled the logistics network not as isolated trips, but as a living Directed Graph to find structural flaws.\n\n"
               "Network Scale:\n"
               "• 1,657 Connected Facilities (Nodes)\n"
               "• 2,783 Logistics Corridors (Edges)\n"
               "• 2,617 Chronic Delay Corridors Identified")
style_content(txBox)

img_path = "outputs/figures/doodle_network.png"
slide.shapes.add_picture(img_path, Inches(5.5), Inches(1.5), width=Inches(4))

# 4. Network Map
slide = prs.slides.add_slide(prs.slide_layouts[5])
set_light_bg(slide)
add_title(slide, "Mapping the Bottlenecks")
img_path = "outputs/figures/network_bottleneck.png"
slide.shapes.add_picture(img_path, Inches(1.5), Inches(1.5), height=Inches(5.5))

# 5. Top 5 Hubs
slide = prs.slides.add_slide(prs.slide_layouts[5])
set_light_bg(slide)
add_title(slide, "The Top 5 Critical Hubs")
img_path = "outputs/figures/top_hubs.png"
slide.shapes.add_picture(img_path, Inches(0.5), Inches(2.0), width=Inches(9))

# 6. ML ETA Models
slide = prs.slides.add_slide(prs.slide_layouts[5])
set_light_bg(slide)
add_title(slide, "Smarter ETAs via Node2Vec")
img_path = "outputs/figures/model_comparison.png"
slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), width=Inches(9))
txBox = slide.shapes.add_textbox(Inches(1), Inches(6.3), Inches(8), Inches(1))
p = txBox.text_frame.paragraphs[0]
p.text = "Result: +8.57% improvement in 15%-Accuracy over OSRM Baseline"
p.font.color.rgb = RGBColor(40, 167, 69)
p.font.bold = True
p.font.size = Pt(22)
p.alignment = PP_ALIGN.CENTER

# 7. FTL Optimization
slide = slide = prs.slides.add_slide(prs.slide_layouts[5])
set_light_bg(slide)
add_title(slide, "FTL vs Carting Framework")
img_path = "outputs/figures/ftl_decision_boundary.png"
slide.shapes.add_picture(img_path, Inches(1.5), Inches(1.5), height=Inches(5.5))

# 8. Business Value
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_light_bg(slide)
add_title(slide, "Strategic Recommendations")
content = slide.placeholders[1]
content.text = ("• Infrastructure Interventions: Upgrade capacity at Gurgaon Bilaspur HB and Bangalore Nelmngla. These are single points of failure causing massive ripple effects.\n"
               "• Financial Impact: Upgrading the top 3 bottleneck hubs to the network median delay recovers an estimated ₹1.05 Crore in revenue at risk.\n"
               "• Product Deployment: Integrate the Graph-Enhanced XGBoost model into the customer-facing tracking app to drastically improve consumer transparency.")
style_content(content)

# 9. Thank You
slide = prs.slides.add_slide(prs.slide_layouts[0])
set_light_bg(slide)
title = slide.shapes.title
title.text = "Thank You"
title.text_frame.paragraphs[0].font.color.rgb = accent_blue
title.text_frame.paragraphs[0].font.size = Pt(54)
subtitle = slide.placeholders[1]
subtitle.text = "Explore the Live Dashboard:\ndelhivery-graph-intelligence-system.streamlit.app"
subtitle.text_frame.paragraphs[0].font.color.rgb = text_dark
subtitle.text_frame.paragraphs[0].font.size = Pt(24)

prs.save("Pitch_Deck.pptx")
print("Presentation saved as Pitch_Deck.pptx")
